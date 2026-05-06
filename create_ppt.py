from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# ── Colors ──────────────────────────────────────────────
BG_DARK   = RGBColor(0x0F, 0x17, 0x2A)   # deep navy
BG_CARD   = RGBColor(0x1A, 0x25, 0x3C)   # card background
ACCENT    = RGBColor(0x00, 0xD4, 0xFF)    # cyan accent
ACCENT2   = RGBColor(0x7C, 0x3A, 0xED)    # purple accent
GREEN     = RGBColor(0x10, 0xB9, 0x81)    # green
ORANGE    = RGBColor(0xF5, 0x9E, 0x0B)    # orange/amber
RED       = RGBColor(0xEF, 0x44, 0x44)    # red
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xCB, 0xD5, 0xE1)    # light gray text
DIMMED    = RGBColor(0x94, 0xA3, 0xB8)    # dimmed text
YELLOW    = RGBColor(0xFA, 0xCC, 0x15)    # yellow

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


# ── Helpers ─────────────────────────────────────────────
def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, shape=MSO_SHAPE.ROUNDED_RECTANGLE, border_color=None, border_width=None):
    shp = slide.shapes.add_shape(shape, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if border_color:
        shp.line.color.rgb = border_color
        shp.line.width = border_width or Pt(2)
    else:
        shp.line.fill.background()
    return shp

def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_multiline(slide, left, top, width, height, lines, default_size=18, default_color=WHITE, align=PP_ALIGN.LEFT, spacing=1.2):
    """lines = list of (text, size, color, bold)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        txt = line_data[0]
        sz  = line_data[1] if len(line_data) > 1 else default_size
        clr = line_data[2] if len(line_data) > 2 else default_color
        bld = line_data[3] if len(line_data) > 3 else False
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(sz)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.font.name = "Segoe UI"
        p.alignment = align
        p.space_after = Pt(sz * (spacing - 1) + 4)
    return txBox

def add_arrow(slide, start_left, start_top, end_left, end_top, color=ACCENT, width=Pt(3)):
    connector = slide.shapes.add_connector(
        1,  # straight connector
        start_left, start_top,
        end_left, end_top
    )
    connector.line.color.rgb = color
    connector.line.width = width
    return connector

def add_circle(slide, left, top, size, fill_color, text="", text_size=24, text_color=WHITE):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    shp.line.fill.background()
    if text:
        tf = shp.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(text_size)
        p.font.color.rgb = text_color
        p.font.bold = True
        p.font.name = "Segoe UI"
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(0)
    return shp

def add_icon_card(slide, left, top, width, height, icon_text, title, desc, icon_color=ACCENT, bg=BG_CARD):
    card = add_shape(slide, left, top, width, height, bg, border_color=RGBColor(0x33, 0x40, 0x55), border_width=Pt(1))
    # icon circle
    circle_size = Inches(0.7)
    add_circle(slide, left + Inches(0.3), top + Inches(0.35), circle_size, icon_color, icon_text, text_size=22)
    # title
    add_text(slide, left + Inches(1.2), top + Inches(0.25), width - Inches(1.5), Inches(0.5), title, size=20, color=WHITE, bold=True)
    # desc
    add_text(slide, left + Inches(1.2), top + Inches(0.75), width - Inches(1.5), height - Inches(1), desc, size=14, color=LIGHT)
    return card

def add_step_number(slide, left, top, number, color=ACCENT):
    size = Inches(0.55)
    add_circle(slide, left, top, size, color, str(number), text_size=20)


# ════════════════════════════════════════════════════════
#  SLIDE 1: TITLE
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

# top decorative bar
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

# badge
add_shape(slide, Inches(4.5), Inches(1.2), Inches(4.3), Inches(0.45), ACCENT2, border_color=None)
add_text(slide, Inches(4.5), Inches(1.22), Inches(4.3), Inches(0.45), "DIGITAL MARKETING WORKSHOP", size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Main title
add_multiline(slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(2.5), [
    ("Smart Email Automation", 48, WHITE, True),
    ("Using AI + n8n", 48, ACCENT, True),
], align=PP_ALIGN.CENTER)

# Subtitle
add_text(slide, Inches(2), Inches(4.3), Inches(9.3), Inches(0.8),
    "Let a Robot Read, Sort & Reply to Your Emails  --  100% FREE",
    size=22, color=LIGHT, align=PP_ALIGN.CENTER)

# bottom cards row
tools = [
    ("n8n", "Automation Builder", ACCENT),
    ("AI", "Ollama (Local AI)", ACCENT2),
    ("Gmail","Email Service", GREEN),
]
card_w = Inches(3.2)
card_h = Inches(1.1)
start_x = Inches(1.85)
for i, (icon, label, clr) in enumerate(tools):
    x = start_x + i * (card_w + Inches(0.4))
    y = Inches(5.6)
    add_shape(slide, x, y, card_w, card_h, BG_CARD, border_color=clr, border_width=Pt(2))
    add_circle(slide, x + Inches(0.2), y + Inches(0.2), Inches(0.65), clr, icon, text_size=16)
    add_text(slide, x + Inches(1.0), y + Inches(0.25), card_w - Inches(1.2), Inches(0.6), label, size=17, color=WHITE, bold=True)

# Footer
add_text(slide, Inches(0), Inches(7.0), W, Inches(0.4),
    "All tools are FREE  |  No coding required  |  Runs on YOUR laptop",
    size=13, color=DIMMED, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
#  SLIDE 2: THE PROBLEM
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), RED)

# Section badge
add_shape(slide, Inches(0.6), Inches(0.4), Inches(2.5), Inches(0.4), RED)
add_text(slide, Inches(0.6), Inches(0.4), Inches(2.5), Inches(0.4), "  THE PROBLEM", size=13, color=WHITE, bold=True)

add_text(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.8),
    "Your Inbox is a MESS", size=40, color=WHITE, bold=True)

add_text(slide, Inches(0.6), Inches(1.7), Inches(10), Inches(0.6),
    "Every day, businesses waste hours sorting emails manually...", size=18, color=LIGHT)

# Email examples
emails = [
    ('"I want to buy 50 red t-shirts"', "ORDER", GREEN, "This is a sale!"),
    ('"What sizes do you have?"', "QUERY", ORANGE, "This is a question"),
    ('"Send me price list for bulk"', "QUERY", ORANGE, "This is a question"),
    ('"Order 10 boxes, deliver to Lahore"', "ORDER", GREEN, "This is a sale!"),
    ('"Do you deliver to Karachi?"', "QUERY", ORANGE, "This is a question"),
]

y = Inches(2.5)
for i, (email_text, etype, clr, note) in enumerate(emails):
    row_h = Inches(0.75)
    # email card
    add_shape(slide, Inches(0.8), y, Inches(7.5), row_h, BG_CARD, border_color=RGBColor(0x33, 0x40, 0x55))
    add_text(slide, Inches(1.0), y + Inches(0.08), Inches(7), row_h, email_text, size=16, color=WHITE)

    # arrow
    add_text(slide, Inches(8.5), y + Inches(0.08), Inches(0.8), row_h, "-->", size=18, color=DIMMED)

    # badge
    badge_w = Inches(1.2)
    add_shape(slide, Inches(9.4), y + Inches(0.12), badge_w, Inches(0.5), clr)
    add_text(slide, Inches(9.4), y + Inches(0.12), badge_w, Inches(0.5), etype, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # note
    add_text(slide, Inches(10.8), y + Inches(0.08), Inches(2.3), row_h, note, size=13, color=DIMMED)

    y += row_h + Inches(0.12)

# Bottom callout
add_shape(slide, Inches(1.5), Inches(6.4), Inches(10.3), Inches(0.7), RGBColor(0x3B, 0x0F, 0x0F), border_color=RED, border_width=Pt(1.5))
add_text(slide, Inches(1.5), Inches(6.45), Inches(10.3), Inches(0.6),
    "Someone reads EVERY email, decides, forwards, and replies. This takes HOURS every day!",
    size=16, color=RGBColor(0xFC, 0xA5, 0xA5), bold=False, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
#  SLIDE 3: THE SOLUTION
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), GREEN)

add_shape(slide, Inches(0.6), Inches(0.4), Inches(2.5), Inches(0.4), GREEN)
add_text(slide, Inches(0.6), Inches(0.4), Inches(2.5), Inches(0.4), "  THE SOLUTION", size=13, color=WHITE, bold=True)

add_text(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.8),
    "Let a ROBOT Do It For You!", size=40, color=WHITE, bold=True)

add_text(slide, Inches(0.6), Inches(1.7), Inches(10), Inches(0.6),
    "An AI-powered email robot that works 24/7 without getting tired", size=18, color=LIGHT)

# Before vs After
# BEFORE card
bx = Inches(0.8)
add_shape(slide, bx, Inches(2.6), Inches(5.5), Inches(4.2), BG_CARD, border_color=RED, border_width=Pt(2))
add_shape(slide, bx + Inches(1.5), Inches(2.35), Inches(2.5), Inches(0.45), RED)
add_text(slide, bx + Inches(1.5), Inches(2.35), Inches(2.5), Inches(0.45), "BEFORE (Manual)", size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

before_items = [
    "Read every email one by one",
    "Decide: Is this order or question?",
    "Copy email, open new window",
    "Forward to correct department",
    "Type a reply to customer",
    "Repeat 100+ times daily",
    "Miss emails, make mistakes",
]
for i, item in enumerate(before_items):
    add_text(slide, bx + Inches(0.3), Inches(3.0) + i * Inches(0.48), Inches(5), Inches(0.45),
        f"  x   {item}", size=15, color=RGBColor(0xFC, 0xA5, 0xA5))

# AFTER card
ax = Inches(7.0)
add_shape(slide, ax, Inches(2.6), Inches(5.5), Inches(4.2), BG_CARD, border_color=GREEN, border_width=Pt(2))
add_shape(slide, ax + Inches(1.5), Inches(2.35), Inches(2.5), Inches(0.45), GREEN)
add_text(slide, ax + Inches(1.5), Inches(2.35), Inches(2.5), Inches(0.45), "AFTER (Robot)", size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

after_items = [
    "Robot reads all emails instantly",
    "AI decides order vs query",
    "Auto-forwards to right team",
    "Auto-replies to customer",
    "Works 24/7, never tired",
    "Handles 1000+ emails/day",
    "Zero mistakes, zero delay",
]
for i, item in enumerate(after_items):
    add_text(slide, ax + Inches(0.3), Inches(3.0) + i * Inches(0.48), Inches(5), Inches(0.45),
        f"  +   {item}", size=15, color=RGBColor(0x6E, 0xE7, 0xB7))


# ════════════════════════════════════════════════════════
#  SLIDE 4: WHAT IS AUTOMATION
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

add_shape(slide, Inches(0.6), Inches(0.4), Inches(2.5), Inches(0.4), ACCENT)
add_text(slide, Inches(0.6), Inches(0.4), Inches(2.5), Inches(0.4), "  BASIC CONCEPT", size=13, color=WHITE, bold=True)

add_text(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.8),
    "What is Automation?", size=40, color=WHITE, bold=True)

# Simple analogy
add_text(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(0.7),
    "Think of it like a WASHING MACHINE for your emails!", size=24, color=ACCENT, bold=True)

# Analogy cards
cards_data = [
    ("Washing Machine", "You put dirty clothes in\nMachine washes automatically\nYou get clean clothes out", ACCENT, "WM"),
    ("Email Robot", "Emails come in\nRobot sorts automatically\nRight people get right emails", GREEN, "ER"),
]
for i, (title, desc, clr, icon) in enumerate(cards_data):
    cx = Inches(1.5) + i * Inches(5.8)
    cy = Inches(3.0)
    add_shape(slide, cx, cy, Inches(5), Inches(3.0), BG_CARD, border_color=clr, border_width=Pt(2))
    add_circle(slide, cx + Inches(1.9), cy + Inches(0.3), Inches(1.2), clr, icon, text_size=26)
    add_text(slide, cx + Inches(0.3), cy + Inches(1.6), Inches(4.4), Inches(0.5), title, size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, cx + Inches(0.3), cy + Inches(2.1), Inches(4.4), Inches(1.0), desc, size=15, color=LIGHT, align=PP_ALIGN.CENTER)

# Key message
add_shape(slide, Inches(2.5), Inches(6.4), Inches(8.3), Inches(0.65), RGBColor(0x0C, 0x2D, 0x1E), border_color=GREEN, border_width=Pt(1.5))
add_text(slide, Inches(2.5), Inches(6.45), Inches(8.3), Inches(0.6),
    "You set it up ONCE, and it works FOREVER automatically!",
    size=17, color=RGBColor(0x6E, 0xE7, 0xB7), bold=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
#  SLIDE 5: OUR TOOLS
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT2)

add_shape(slide, Inches(0.6), Inches(0.4), Inches(2.5), Inches(0.4), ACCENT2)
add_text(slide, Inches(0.6), Inches(0.4), Inches(2.5), Inches(0.4), "  OUR TOOLKIT", size=13, color=WHITE, bold=True)

add_text(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.8),
    "3 Free Tools = 1 Powerful Robot", size=40, color=WHITE, bold=True)

add_text(slide, Inches(0.6), Inches(1.7), Inches(10), Inches(0.6),
    "Everything runs on YOUR laptop. No subscriptions. No credit cards.", size=18, color=LIGHT)

# Tool cards - bigger and more descriptive
tools_info = [
    ("n8n", "The Robot Builder", "Like LEGO blocks for automation.\nYou connect blocks together.\nEach block does one job.\nNo coding needed!", ACCENT, "FREE"),
    ("Ollama", "The AI Brain", "Like ChatGPT on your laptop.\nReads emails & understands them.\nDecides: Order or Query?\nWorks without internet!", ACCENT2, "FREE"),
    ("Gmail", "The Email Gate", "Your regular email account.\nRobot watches for new emails.\nSends auto-replies.\nForwards to departments.", GREEN, "FREE"),
]

for i, (name, subtitle, desc, clr, cost) in enumerate(tools_info):
    cx = Inches(0.6) + i * Inches(4.15)
    cy = Inches(2.6)
    cw = Inches(3.85)
    ch = Inches(4.3)
    add_shape(slide, cx, cy, cw, ch, BG_CARD, border_color=clr, border_width=Pt(2))

    # icon circle
    add_circle(slide, cx + Inches(1.3), cy + Inches(0.3), Inches(1.2), clr, name, text_size=22)

    # subtitle
    add_text(slide, cx + Inches(0.2), cy + Inches(1.65), cw - Inches(0.4), Inches(0.4), subtitle, size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # description
    add_text(slide, cx + Inches(0.3), cy + Inches(2.15), cw - Inches(0.6), Inches(1.5), desc, size=14, color=LIGHT, align=PP_ALIGN.CENTER)

    # cost badge
    add_shape(slide, cx + Inches(1.2), cy + Inches(3.7), Inches(1.4), Inches(0.4), clr)
    add_text(slide, cx + Inches(1.2), cy + Inches(3.7), Inches(1.4), Inches(0.4), cost, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
#  SLIDE 6: HOW IT WORKS - THE FLOW
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

add_shape(slide, Inches(0.6), Inches(0.4), Inches(2.5), Inches(0.4), ACCENT)
add_text(slide, Inches(0.6), Inches(0.4), Inches(2.5), Inches(0.4), "  THE FLOW", size=13, color=WHITE, bold=True)

add_text(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.8),
    "How Our Email Robot Works", size=40, color=WHITE, bold=True)

# Flow diagram using shapes
# Row 1: Email comes in
step_y1 = Inches(2.3)
step_h = Inches(1.0)

# Step 1: Customer Email
add_shape(slide, Inches(0.5), step_y1, Inches(2.8), step_h, BG_CARD, border_color=ACCENT, border_width=Pt(2))
add_step_number(slide, Inches(0.6), step_y1 - Inches(0.15), 1, ACCENT)
add_text(slide, Inches(0.7), step_y1 + Inches(0.1), Inches(2.4), Inches(0.35), "Customer Sends", size=14, color=DIMMED, align=PP_ALIGN.CENTER)
add_text(slide, Inches(0.7), step_y1 + Inches(0.45), Inches(2.4), Inches(0.4), "EMAIL", size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Arrow
add_text(slide, Inches(3.4), step_y1 + Inches(0.2), Inches(0.7), step_h, ">>>", size=24, color=ACCENT, bold=True)

# Step 2: Robot Gets It
add_shape(slide, Inches(4.0), step_y1, Inches(2.8), step_h, BG_CARD, border_color=ACCENT, border_width=Pt(2))
add_step_number(slide, Inches(4.1), step_y1 - Inches(0.15), 2, ACCENT)
add_text(slide, Inches(4.2), step_y1 + Inches(0.1), Inches(2.4), Inches(0.35), "n8n Robot", size=14, color=DIMMED, align=PP_ALIGN.CENTER)
add_text(slide, Inches(4.2), step_y1 + Inches(0.45), Inches(2.4), Inches(0.4), "CATCHES IT", size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Arrow
add_text(slide, Inches(6.9), step_y1 + Inches(0.2), Inches(0.7), step_h, ">>>", size=24, color=ACCENT, bold=True)

# Step 3: AI Reads
add_shape(slide, Inches(7.5), step_y1, Inches(2.8), step_h, BG_CARD, border_color=ACCENT2, border_width=Pt(2))
add_step_number(slide, Inches(7.6), step_y1 - Inches(0.15), 3, ACCENT2)
add_text(slide, Inches(7.7), step_y1 + Inches(0.1), Inches(2.4), Inches(0.35), "AI (Ollama)", size=14, color=DIMMED, align=PP_ALIGN.CENTER)
add_text(slide, Inches(7.7), step_y1 + Inches(0.45), Inches(2.4), Inches(0.4), "READS EMAIL", size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Arrow
add_text(slide, Inches(10.4), step_y1 + Inches(0.2), Inches(0.7), step_h, ">>>", size=24, color=ACCENT2, bold=True)

# Step 4: AI Decides
add_shape(slide, Inches(11.0), step_y1, Inches(1.8), step_h, BG_CARD, border_color=YELLOW, border_width=Pt(2))
add_step_number(slide, Inches(11.1), step_y1 - Inches(0.15), 4, YELLOW)
add_text(slide, Inches(11.0), step_y1 + Inches(0.1), Inches(1.8), Inches(0.35), "AI Decides", size=14, color=DIMMED, align=PP_ALIGN.CENTER)
add_text(slide, Inches(11.0), step_y1 + Inches(0.45), Inches(1.8), Inches(0.4), "ORDER?", size=20, color=YELLOW, bold=True, align=PP_ALIGN.CENTER)

# Row 2 - Split into two paths
# ORDER PATH (left)
step_y2 = Inches(4.2)

add_shape(slide, Inches(1.0), step_y2, Inches(5.0), Inches(2.6), BG_CARD, border_color=GREEN, border_width=Pt(2))
add_shape(slide, Inches(2.5), step_y2 - Inches(0.2), Inches(2.0), Inches(0.4), GREEN)
add_text(slide, Inches(2.5), step_y2 - Inches(0.2), Inches(2.0), Inches(0.4), "YES = ORDER", size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_multiline(slide, Inches(1.4), step_y2 + Inches(0.4), Inches(4.2), Inches(2.0), [
    ("Step 5a:", 14, DIMMED, False),
    ("Forward to Order Department", 18, WHITE, True),
    ("", 8, WHITE, False),
    ("Step 6a:", 14, DIMMED, False),
    ("Reply: 'Thank you for your order!'", 16, RGBColor(0x6E, 0xE7, 0xB7), False),
])

# QUERY PATH (right)
add_shape(slide, Inches(7.0), step_y2, Inches(5.5), Inches(2.6), BG_CARD, border_color=ORANGE, border_width=Pt(2))
add_shape(slide, Inches(8.7), step_y2 - Inches(0.2), Inches(2.0), Inches(0.4), ORANGE)
add_text(slide, Inches(8.7), step_y2 - Inches(0.2), Inches(2.0), Inches(0.4), "NO = QUERY", size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_multiline(slide, Inches(7.4), step_y2 + Inches(0.4), Inches(4.8), Inches(2.0), [
    ("Step 5b:", 14, DIMMED, False),
    ("Forward to Customer Support", 18, WHITE, True),
    ("", 8, WHITE, False),
    ("Step 6b:", 14, DIMMED, False),
    ("Reply: 'We received your query!'", 16, RGBColor(0xFD, 0xE6, 0x8A), False),
])


# ════════════════════════════════════════════════════════
#  SLIDE 7: WHAT IS n8n (Visual)
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

add_shape(slide, Inches(0.6), Inches(0.4), Inches(2.5), Inches(0.4), ACCENT)
add_text(slide, Inches(0.6), Inches(0.4), Inches(2.5), Inches(0.4), "  TOOL #1", size=13, color=WHITE, bold=True)

add_text(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.8),
    "n8n = Your Robot Factory", size=40, color=WHITE, bold=True)

add_text(slide, Inches(0.6), Inches(1.7), Inches(10), Inches(0.6),
    "You build robots by connecting blocks. Like LEGO. No coding.", size=20, color=LIGHT)

# Visual representation of n8n blocks
block_y = Inches(2.8)
blocks = [
    ("Trigger", "When email\narrives", ACCENT, Inches(0.5)),
    ("AI", "Read email\nwith AI", ACCENT2, Inches(3.3)),
    ("IF", "Order or\nQuery?", YELLOW, Inches(6.1)),
    ("Email", "Forward &\nReply", GREEN, Inches(8.9)),
]

for name, desc, clr, bx in blocks:
    add_shape(slide, bx, block_y, Inches(2.3), Inches(1.6), BG_CARD, border_color=clr, border_width=Pt(2.5))
    add_text(slide, bx, block_y + Inches(0.15), Inches(2.3), Inches(0.35), name, size=18, color=clr, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, bx, block_y + Inches(0.55), Inches(2.3), Inches(0.9), desc, size=14, color=LIGHT, align=PP_ALIGN.CENTER)

# Arrows between blocks
for i in range(3):
    ax = [Inches(2.85), Inches(5.65), Inches(8.45)][i]
    add_text(slide, ax, block_y + Inches(0.4), Inches(0.5), Inches(0.6), ">>", size=22, color=ACCENT, bold=True)

# Key points
pts = [
    ("Drag & Drop", "Just drag blocks onto the screen and connect them"),
    ("No Coding", "Everything is visual - click, fill forms, done"),
    ("Free Forever", "Open source, runs on your laptop, no limits"),
    ("Works Everywhere", "Same blocks on laptop or cloud - learn once, use anywhere"),
]
py = Inches(4.8)
for i, (title, desc) in enumerate(pts):
    px = Inches(0.6) + (i % 2) * Inches(6.3)
    row_y = py + (i // 2) * Inches(1.0)
    add_circle(slide, px, row_y + Inches(0.05), Inches(0.45), ACCENT, str(i+1), text_size=16)
    add_text(slide, px + Inches(0.6), row_y, Inches(5.2), Inches(0.35), title, size=16, color=WHITE, bold=True)
    add_text(slide, px + Inches(0.6), row_y + Inches(0.35), Inches(5.2), Inches(0.4), desc, size=13, color=LIGHT)


# ════════════════════════════════════════════════════════
#  SLIDE 8: WHAT IS OLLAMA
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT2)

add_shape(slide, Inches(0.6), Inches(0.4), Inches(2.5), Inches(0.4), ACCENT2)
add_text(slide, Inches(0.6), Inches(0.4), Inches(2.5), Inches(0.4), "  TOOL #2", size=13, color=WHITE, bold=True)

add_text(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.8),
    "Ollama = ChatGPT on Your Laptop", size=40, color=WHITE, bold=True)

add_text(slide, Inches(0.6), Inches(1.7), Inches(10), Inches(0.6),
    "Same AI power, but it runs on YOUR computer. Free. Private. No internet needed.", size=20, color=LIGHT)

# Comparison: ChatGPT vs Ollama
headers = [("ChatGPT", ORANGE, Inches(1.0)), ("Ollama", GREEN, Inches(7.2))]
for title, clr, hx in headers:
    add_shape(slide, hx, Inches(2.6), Inches(5.0), Inches(0.55), clr)
    add_text(slide, hx, Inches(2.6), Inches(5.0), Inches(0.55), title, size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

comparisons = [
    ("Runs on OpenAI's servers", "Runs on YOUR laptop"),
    ("Needs internet always", "Works without internet"),
    ("$20/month subscription", "100% FREE forever"),
    ("Your data goes to cloud", "Your data stays private"),
    ("They can shut it off", "You control everything"),
]

for i, (chatgpt, ollama) in enumerate(comparisons):
    ry = Inches(3.35) + i * Inches(0.65)
    # ChatGPT column
    add_shape(slide, Inches(1.0), ry, Inches(5.0), Inches(0.55), BG_CARD)
    add_text(slide, Inches(1.2), ry + Inches(0.05), Inches(4.6), Inches(0.45), f"  x  {chatgpt}", size=14, color=RGBColor(0xFC, 0xA5, 0xA5))
    # Ollama column
    add_shape(slide, Inches(7.2), ry, Inches(5.0), Inches(0.55), BG_CARD)
    add_text(slide, Inches(7.4), ry + Inches(0.05), Inches(4.6), Inches(0.45), f"  +  {ollama}", size=14, color=RGBColor(0x6E, 0xE7, 0xB7))

# vs circle
add_circle(slide, Inches(6.0), Inches(3.8), Inches(1.1), ACCENT2, "VS", text_size=22)

# Bottom note
add_shape(slide, Inches(2.5), Inches(6.6), Inches(8.3), Inches(0.55), RGBColor(0x1A, 0x15, 0x3C), border_color=ACCENT2, border_width=Pt(1.5))
add_text(slide, Inches(2.5), Inches(6.63), Inches(8.3), Inches(0.5),
    "Later, if you want more power, just swap Ollama with ChatGPT/Claude API. Same workflow!",
    size=14, color=RGBColor(0xC4, 0xB5, 0xFD), align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
#  SLIDE 9: VOCABULARY - Only 4 Words
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), YELLOW)

add_shape(slide, Inches(0.6), Inches(0.4), Inches(2.5), Inches(0.4), YELLOW)
add_text(slide, Inches(0.6), Inches(0.4), Inches(2.5), Inches(0.4), "  VOCABULARY", size=13, color=BG_DARK, bold=True)

add_text(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.8),
    "Only 4 Words to Remember!", size=40, color=WHITE, bold=True)

add_text(slide, Inches(0.6), Inches(1.7), Inches(10), Inches(0.6),
    "That's it. Just 4 words and you can build any automation.", size=18, color=LIGHT)

vocab = [
    ("NODE", "One block that does ONE job", "Like one worker in a factory", ACCENT, "1"),
    ("WORKFLOW", "All blocks connected together", "The whole factory assembly line", GREEN, "2"),
    ("TRIGGER", "The block that STARTS everything", "The alarm clock that wakes the robot", ORANGE, "3"),
    ("CONNECTION", "The line between two blocks", "The conveyor belt between workers", ACCENT2, "4"),
]

for i, (word, meaning, analogy, clr, num) in enumerate(vocab):
    cx = Inches(0.6) + (i % 2) * Inches(6.3)
    cy = Inches(2.7) + (i // 2) * Inches(2.1)
    cw = Inches(5.9)
    ch = Inches(1.8)

    add_shape(slide, cx, cy, cw, ch, BG_CARD, border_color=clr, border_width=Pt(2))
    add_circle(slide, cx + Inches(0.2), cy + Inches(0.25), Inches(0.5), clr, num, text_size=20)
    add_text(slide, cx + Inches(0.85), cy + Inches(0.2), Inches(3.5), Inches(0.45), word, size=26, color=clr, bold=True)
    add_text(slide, cx + Inches(0.85), cy + Inches(0.7), cw - Inches(1.2), Inches(0.35), meaning, size=15, color=WHITE)
    add_text(slide, cx + Inches(0.85), cy + Inches(1.1), cw - Inches(1.2), Inches(0.45), f'Example: "{analogy}"', size=13, color=DIMMED)


# ════════════════════════════════════════════════════════
#  SLIDE 10: REAL EXAMPLE DEMO
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), GREEN)

add_shape(slide, Inches(0.6), Inches(0.4), Inches(2.5), Inches(0.4), GREEN)
add_text(slide, Inches(0.6), Inches(0.4), Inches(2.5), Inches(0.4), "  LIVE EXAMPLE", size=13, color=WHITE, bold=True)

add_text(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.8),
    "Watch the Robot in Action!", size=40, color=WHITE, bold=True)

# Example 1: ORDER
add_shape(slide, Inches(0.5), Inches(2.2), Inches(6.0), Inches(2.2), BG_CARD, border_color=GREEN, border_width=Pt(2))
add_shape(slide, Inches(0.5), Inches(2.2), Inches(6.0), Inches(0.5), GREEN)
add_text(slide, Inches(0.7), Inches(2.22), Inches(5.5), Inches(0.45), "EXAMPLE 1: Customer sends this email", size=14, color=WHITE, bold=True)

add_multiline(slide, Inches(0.8), Inches(2.85), Inches(5.4), Inches(1.5), [
    ('Subject: Order for t-shirts', 14, LIGHT, False),
    ('', 6, WHITE, False),
    ('"Hi, I want to order 50 red t-shirts.', 15, WHITE, False),
    ('Size L. Deliver to Lahore. COD payment."', 15, WHITE, False),
])

# Arrow
add_text(slide, Inches(6.6), Inches(2.8), Inches(0.8), Inches(1.0), ">>>", size=24, color=GREEN, bold=True)

# Result
add_shape(slide, Inches(7.2), Inches(2.2), Inches(5.6), Inches(2.2), BG_CARD, border_color=GREEN, border_width=Pt(2))
add_shape(slide, Inches(7.2), Inches(2.2), Inches(5.6), Inches(0.5), GREEN)
add_text(slide, Inches(7.4), Inches(2.22), Inches(5.2), Inches(0.45), "ROBOT DOES THIS:", size=14, color=WHITE, bold=True)

add_multiline(slide, Inches(7.5), Inches(2.85), Inches(5.1), Inches(1.5), [
    ('AI says: "ORDER"', 15, RGBColor(0x6E, 0xE7, 0xB7), True),
    ('', 6, WHITE, False),
    ('+ Forwards to orders@company.com', 14, WHITE, False),
    ('+ Replies: "Thank you for your order!"', 14, WHITE, False),
])

# Example 2: QUERY
add_shape(slide, Inches(0.5), Inches(4.8), Inches(6.0), Inches(2.2), BG_CARD, border_color=ORANGE, border_width=Pt(2))
add_shape(slide, Inches(0.5), Inches(4.8), Inches(6.0), Inches(0.5), ORANGE)
add_text(slide, Inches(0.7), Inches(4.82), Inches(5.5), Inches(0.45), "EXAMPLE 2: Customer sends this email", size=14, color=WHITE, bold=True)

add_multiline(slide, Inches(0.8), Inches(5.45), Inches(5.4), Inches(1.5), [
    ('Subject: Question about delivery', 14, LIGHT, False),
    ('', 6, WHITE, False),
    ('"Hello, do you deliver to Karachi?', 15, WHITE, False),
    ('What are your delivery charges?"', 15, WHITE, False),
])

# Arrow
add_text(slide, Inches(6.6), Inches(5.4), Inches(0.8), Inches(1.0), ">>>", size=24, color=ORANGE, bold=True)

# Result
add_shape(slide, Inches(7.2), Inches(4.8), Inches(5.6), Inches(2.2), BG_CARD, border_color=ORANGE, border_width=Pt(2))
add_shape(slide, Inches(7.2), Inches(4.8), Inches(5.6), Inches(0.5), ORANGE)
add_text(slide, Inches(7.4), Inches(4.82), Inches(5.2), Inches(0.45), "ROBOT DOES THIS:", size=14, color=WHITE, bold=True)

add_multiline(slide, Inches(7.5), Inches(5.45), Inches(5.1), Inches(1.5), [
    ('AI says: "QUERY"', 15, RGBColor(0xFD, 0xE6, 0x8A), True),
    ('', 6, WHITE, False),
    ('+ Forwards to support@company.com', 14, WHITE, False),
    ('+ Replies: "We received your query!"', 14, WHITE, False),
])


# ════════════════════════════════════════════════════════
#  SLIDE 11: INSTALLATION STEPS
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

add_shape(slide, Inches(0.6), Inches(0.4), Inches(2.5), Inches(0.4), ACCENT)
add_text(slide, Inches(0.6), Inches(0.4), Inches(2.5), Inches(0.4), "  INSTALLATION", size=13, color=WHITE, bold=True)

add_text(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.8),
    "Setting Up (One Time Only)", size=40, color=WHITE, bold=True)

# Step cards
steps = [
    ("1", "Install Ollama", "Go to ollama.com\nClick Download > Windows\nInstall (Next, Next, Done)", ACCENT2,
     "Then open Terminal and type:\nollama pull llama3.2"),
    ("2", "Install Node.js", "Go to nodejs.org\nDownload LTS (green button)\nInstall (Next, Next, Done)", ACCENT,
     "This is needed to run n8n\non your laptop"),
    ("3", "Install n8n", "Open Terminal and type:\nnpm install -g n8n\nWait for it to finish", GREEN,
     "Then type:  n8n start\nOpen browser: localhost:5678"),
]

for i, (num, title, desc, clr, note) in enumerate(steps):
    cx = Inches(0.4) + i * Inches(4.25)
    cy = Inches(2.2)
    cw = Inches(4.0)
    ch = Inches(4.8)

    add_shape(slide, cx, cy, cw, ch, BG_CARD, border_color=clr, border_width=Pt(2))
    add_circle(slide, cx + Inches(1.5), cy + Inches(0.3), Inches(0.9), clr, num, text_size=32)
    add_text(slide, cx + Inches(0.2), cy + Inches(1.35), cw - Inches(0.4), Inches(0.45), title, size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, cx + Inches(0.3), cy + Inches(1.9), cw - Inches(0.6), Inches(1.5), desc, size=14, color=LIGHT, align=PP_ALIGN.CENTER)

    # Note box
    add_shape(slide, cx + Inches(0.2), cy + Inches(3.4), cw - Inches(0.4), Inches(1.1), RGBColor(0x0F, 0x17, 0x2A), border_color=clr, border_width=Pt(1))
    add_text(slide, cx + Inches(0.35), cy + Inches(3.5), cw - Inches(0.7), Inches(0.9), note, size=12, color=DIMMED, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
#  SLIDE 12: DAILY STARTUP
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), GREEN)

add_shape(slide, Inches(0.6), Inches(0.4), Inches(2.5), Inches(0.4), GREEN)
add_text(slide, Inches(0.6), Inches(0.4), Inches(2.5), Inches(0.4), "  DAILY USE", size=13, color=WHITE, bold=True)

add_text(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.8),
    "Start Your Robot Every Day", size=40, color=WHITE, bold=True)

add_text(slide, Inches(0.6), Inches(1.7), Inches(10), Inches(0.6),
    "Just 3 simple commands every morning. That's it!", size=20, color=LIGHT)

daily_steps = [
    ("1", "Start AI Brain", "Open Terminal, type:", "ollama serve", ACCENT2),
    ("2", "Start Robot", "Open NEW Terminal, type:", "n8n start", ACCENT),
    ("3", "Check Robot", "Open Browser, go to:", "localhost:5678", GREEN),
]

for i, (num, title, instruction, command, clr) in enumerate(daily_steps):
    cx = Inches(0.6) + i * Inches(4.2)
    cy = Inches(2.8)
    cw = Inches(3.9)
    ch = Inches(3.5)

    add_shape(slide, cx, cy, cw, ch, BG_CARD, border_color=clr, border_width=Pt(2))
    add_circle(slide, cx + Inches(1.4), cy + Inches(0.3), Inches(1.0), clr, num, text_size=36)
    add_text(slide, cx + Inches(0.2), cy + Inches(1.5), cw - Inches(0.4), Inches(0.45), title, size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, cx + Inches(0.2), cy + Inches(2.0), cw - Inches(0.4), Inches(0.35), instruction, size=14, color=DIMMED, align=PP_ALIGN.CENTER)
    # Command box
    add_shape(slide, cx + Inches(0.3), cy + Inches(2.5), cw - Inches(0.6), Inches(0.6), RGBColor(0x0A, 0x0F, 0x1A), border_color=clr, border_width=Pt(1.5))
    add_text(slide, cx + Inches(0.3), cy + Inches(2.55), cw - Inches(0.6), Inches(0.55), command, size=18, color=clr, bold=True, align=PP_ALIGN.CENTER, font_name="Consolas")

# Make sure Active note
add_shape(slide, Inches(2.5), Inches(6.6), Inches(8.3), Inches(0.55), RGBColor(0x0C, 0x2D, 0x1E), border_color=GREEN, border_width=Pt(1.5))
add_text(slide, Inches(2.5), Inches(6.63), Inches(8.3), Inches(0.5),
    "Make sure your workflow toggle is set to ACTIVE (green) in n8n!",
    size=15, color=RGBColor(0x6E, 0xE7, 0xB7), bold=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
#  SLIDE 13: FUTURE UPGRADES
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT2)

add_shape(slide, Inches(0.6), Inches(0.4), Inches(2.5), Inches(0.4), ACCENT2)
add_text(slide, Inches(0.6), Inches(0.4), Inches(2.5), Inches(0.4), "  FUTURE", size=13, color=WHITE, bold=True)

add_text(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.8),
    "When You're Ready to Upgrade", size=40, color=WHITE, bold=True)

add_text(slide, Inches(0.6), Inches(1.7), Inches(10), Inches(0.6),
    "What you learn today works EVERYWHERE. No need to relearn!", size=20, color=LIGHT)

upgrades = [
    ("TODAY", "Free + Local", "Ollama on laptop\nn8n on laptop\nFree forever\nLearning & testing", ACCENT, "NOW"),
    ("LEVEL 2", "Better AI", "Swap Ollama with\nChatGPT or Claude API\nMore accurate results\nJust change 1 block!", ACCENT2, "$$$"),
    ("LEVEL 3", "Cloud 24/7", "Move n8n to cloud\nRuns without laptop\nn8n.io subscription\nSame workflow!", GREEN, "$$$"),
]

for i, (level, subtitle, desc, clr, badge) in enumerate(upgrades):
    cx = Inches(0.6) + i * Inches(4.2)
    cy = Inches(2.6)
    cw = Inches(3.9)
    ch = Inches(4.0)

    add_shape(slide, cx, cy, cw, ch, BG_CARD, border_color=clr, border_width=Pt(2))
    add_shape(slide, cx + Inches(0.2), cy + Inches(0.2), Inches(1.5), Inches(0.4), clr)
    add_text(slide, cx + Inches(0.2), cy + Inches(0.2), Inches(1.5), Inches(0.4), level, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_text(slide, cx + Inches(0.2), cy + Inches(0.8), cw - Inches(0.4), Inches(0.4), subtitle, size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, cx + Inches(0.3), cy + Inches(1.4), cw - Inches(0.6), Inches(2.0), desc, size=15, color=LIGHT, align=PP_ALIGN.CENTER)

    # Badge
    add_circle(slide, cx + cw - Inches(1.0), cy + Inches(0.1), Inches(0.65), clr, badge, text_size=13)

# Arrow between levels
add_text(slide, Inches(4.55), Inches(4.2), Inches(0.7), Inches(0.6), ">>>", size=22, color=DIMMED, bold=True)
add_text(slide, Inches(8.75), Inches(4.2), Inches(0.7), Inches(0.6), ">>>", size=22, color=DIMMED, bold=True)

# KEY point
add_shape(slide, Inches(1.5), Inches(6.9), Inches(10.3), Inches(0.45), RGBColor(0x1A, 0x15, 0x3C), border_color=ACCENT2, border_width=Pt(1.5))
add_text(slide, Inches(1.5), Inches(6.9), Inches(10.3), Inches(0.45),
    "Your workflow is the SAME at every level. Learn once, upgrade anytime!",
    size=15, color=RGBColor(0xC4, 0xB5, 0xFD), bold=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
#  SLIDE 14: MORE AUTOMATION IDEAS
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ORANGE)

add_shape(slide, Inches(0.6), Inches(0.4), Inches(2.5), Inches(0.4), ORANGE)
add_text(slide, Inches(0.6), Inches(0.4), Inches(2.5), Inches(0.4), "  MORE IDEAS", size=13, color=WHITE, bold=True)

add_text(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.8),
    "What Else Can You Automate?", size=40, color=WHITE, bold=True)

add_text(slide, Inches(0.6), Inches(1.7), Inches(10), Inches(0.6),
    "Same n8n blocks, different workflows. Sky is the limit!", size=18, color=LIGHT)

ideas = [
    ("DM", "Auto-Reply to\nInstagram DMs", ACCENT),
    ("LEAD", "Score Leads\nHot / Warm / Cold", RED),
    ("POST", "Auto-Generate\nSocial Media Posts", ACCENT2),
    ("CRM", "Auto-Add Leads\nto Google Sheets", GREEN),
    ("WA", "WhatsApp\nAuto-Responder", RGBColor(0x25, 0xD3, 0x66)),
    ("INV", "Auto-Generate\nInvoices", ORANGE),
    ("REV", "Respond to\nGoogle Reviews", YELLOW),
    ("SPY", "Monitor\nCompetitors Daily", RGBColor(0xEC, 0x48, 0x99)),
]

for i, (icon, label, clr) in enumerate(ideas):
    col = i % 4
    row = i // 4
    cx = Inches(0.5) + col * Inches(3.2)
    cy = Inches(2.6) + row * Inches(2.2)
    cw = Inches(2.9)
    ch = Inches(1.9)

    add_shape(slide, cx, cy, cw, ch, BG_CARD, border_color=clr, border_width=Pt(2))
    add_circle(slide, cx + Inches(0.9), cy + Inches(0.2), Inches(1.0), clr, icon, text_size=18)
    add_text(slide, cx + Inches(0.1), cy + Inches(1.25), cw - Inches(0.2), Inches(0.65), label, size=14, color=WHITE, bold=False, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
#  SLIDE 15: KEY TAKEAWAY
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

# Big quote
add_text(slide, Inches(1.5), Inches(1.5), Inches(10.3), Inches(1.5),
    "KEY TAKEAWAY", size=20, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

add_multiline(slide, Inches(1.0), Inches(2.5), Inches(11.3), Inches(2.5), [
    ('"You don\'t need to be a programmer', 36, WHITE, True),
    ('to automate your business."', 36, WHITE, True),
], align=PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(4.5), Inches(10.3), Inches(0.8),
    "Just think: What do I do again and again? Then let a robot do it.",
    size=22, color=LIGHT, align=PP_ALIGN.CENTER)

# 3 summary points
summary = [
    ("FREE Tools", "n8n + Ollama + Gmail\nNo money needed to start", ACCENT),
    ("NO Coding", "Drag, drop, connect blocks\nFill simple forms", GREEN),
    ("REAL Results", "Save hours every day\nZero mistakes, 24/7 work", ORANGE),
]

for i, (title, desc, clr) in enumerate(summary):
    cx = Inches(1.0) + i * Inches(3.9)
    cy = Inches(5.3)
    add_shape(slide, cx, cy, Inches(3.5), Inches(1.5), BG_CARD, border_color=clr, border_width=Pt(2))
    add_text(slide, cx + Inches(0.2), cy + Inches(0.15), Inches(3.1), Inches(0.4), title, size=20, color=clr, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, cx + Inches(0.2), cy + Inches(0.6), Inches(3.1), Inches(0.8), desc, size=14, color=LIGHT, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
#  SLIDE 16: LET'S BUILD IT!
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), GREEN)

add_text(slide, Inches(1.5), Inches(1.5), Inches(10.3), Inches(1.0),
    "THEORY DONE!", size=24, color=DIMMED, align=PP_ALIGN.CENTER)

add_multiline(slide, Inches(1.0), Inches(2.5), Inches(11.3), Inches(2.0), [
    ("Now Let's BUILD It!", 52, GREEN, True),
    ("", 10, WHITE, False),
    ("Hands-on Practical Session", 28, WHITE, False),
], align=PP_ALIGN.CENTER)

# What we'll do boxes
practical = [
    ("1", "Install all 3 tools on your laptop", ACCENT2),
    ("2", "Build the email robot step by step", ACCENT),
    ("3", "Test it with real emails", GREEN),
    ("4", "Make it run automatically", ORANGE),
]

for i, (num, desc, clr) in enumerate(practical):
    cx = Inches(2.5)
    cy = Inches(5.0) + i * Inches(0.55)
    add_circle(slide, cx, cy, Inches(0.4), clr, num, text_size=16)
    add_text(slide, cx + Inches(0.55), cy + Inches(0.02), Inches(8), Inches(0.4), desc, size=17, color=WHITE)


# ════════════════════════════════════════════════════════
#  SAVE
# ════════════════════════════════════════════════════════
output_path = r"c:\Users\Adnan\Desktop\DM\Email_Automation_Lecture.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
